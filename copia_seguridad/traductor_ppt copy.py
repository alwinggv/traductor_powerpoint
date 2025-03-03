#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
Traductor de Presentaciones PowerPoint
Este script permite traducir automáticamente el contenido de presentaciones PowerPoint
de un idioma a otro utilizando la biblioteca deep-translator.
"""

import os
import sys
import time
import shutil
import re
from pptx import Presentation
from deep_translator import GoogleTranslator
from tqdm import tqdm
from dotenv import load_dotenv

# Cargar variables de entorno si existe un archivo .env
load_dotenv()

# Definir rutas de carpetas
CARPETA_ORIGINALES = "originales"
CARPETA_TRADUCIDOS = "traducidos"
CARPETA_COPIA_SEGURIDAD = "copia_seguridad"  # Añadimos la carpeta de copia de seguridad pero no la tocamos

# Crear carpetas si no existen
os.makedirs(CARPETA_ORIGINALES, exist_ok=True)
os.makedirs(CARPETA_TRADUCIDOS, exist_ok=True)

def traducir_texto(texto, idioma_destino='en', idioma_origen='auto'):
    """
    Traduce un texto al idioma especificado.
    
    Args:
        texto (str): Texto a traducir
        idioma_destino (str): Código del idioma de destino (por defecto 'en' para inglés)
        idioma_origen (str): Código del idioma de origen (por defecto 'auto' para detección automática)
    
    Returns:
        str: Texto traducido o el texto original si hay un error
    """
    if not texto or not texto.strip():
        return texto
    
    try:
        # Añadir un pequeño retraso para evitar bloqueos de la API
        time.sleep(0.5)
        
        # Usar GoogleTranslator de deep_translator
        translator = GoogleTranslator(source=idioma_origen, target=idioma_destino)
        resultado = translator.translate(texto)
        
        # Verificar que el resultado no sea None
        if resultado is None:
            print(f"Advertencia: La traducción de '{texto}' devolvió None. Se usará el texto original.")
            return texto
            
        return resultado
    except Exception as e:
        print(f"Error al traducir el texto '{texto}': {e}")
        return texto

def traducir_presentacion(ruta_archivo, idioma_destino='en', idioma_origen='auto'):
    """
    Traduce el contenido de una presentación PowerPoint.
    
    Args:
        ruta_archivo (str): Ruta al archivo PowerPoint
        idioma_destino (str): Código del idioma de destino
        idioma_origen (str): Código del idioma de origen
    
    Returns:
        str: Ruta del archivo traducido
    """
    try:
        # Cargar la presentación
        presentacion = Presentation(ruta_archivo)
        
        # Obtener el nombre del archivo sin extensión
        nombre_base = os.path.basename(ruta_archivo)
        nombre_sin_extension = os.path.splitext(nombre_base)[0]
        
        # Crear el nombre del archivo traducido
        archivo_traducido = f"{nombre_sin_extension}_traducido_{idioma_destino}.pptx"
        ruta_archivo_traducido = os.path.join(CARPETA_TRADUCIDOS, archivo_traducido)
        
        # Verificar si el archivo ya existe y está en uso
        if os.path.exists(ruta_archivo_traducido):
            try:
                # Intentar abrir el archivo para verificar si está en uso
                with open(ruta_archivo_traducido, 'a'):
                    pass
            except PermissionError:
                # Si está en uso, crear un nombre alternativo con timestamp
                import datetime
                timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
                archivo_traducido = f"{nombre_sin_extension}_traducido_{idioma_destino}_{timestamp}.pptx"
                ruta_archivo_traducido = os.path.join(CARPETA_TRADUCIDOS, archivo_traducido)
                print(f"El archivo original está en uso. Se usará un nombre alternativo: {archivo_traducido}")
        
        # Contar el número total de elementos a traducir para la barra de progreso
        total_elementos = sum(len(slide.shapes) for slide in presentacion.slides)
        
        # Inicializar la barra de progreso
        barra_progreso = tqdm(total=total_elementos, desc="Traduciendo presentación")
        
        # Procesar cada diapositiva
        for slide in presentacion.slides:
            # Procesar cada forma en la diapositiva
            for shape in slide.shapes:
                try:
                    # Verificar si la forma tiene texto
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        # Procesar cada párrafo en el marco de texto
                        for paragraph in shape.text_frame.paragraphs:
                            # Si el párrafo está vacío, continuar con el siguiente
                            if not paragraph.text.strip():
                                continue
                                
                            # Recopilar información sobre cada run (texto y formato)
                            runs_info = []
                            for run in paragraph.runs:
                                if run.text and run.text.strip():
                                    runs_info.append({
                                        'run': run,
                                        'texto': run.text
                                    })
                            
                            # Si no hay runs con texto, continuar con el siguiente párrafo
                            if not runs_info:
                                continue
                                
                            # Extraer solo el texto de cada run para formar el párrafo completo
                            texto_completo = ''.join(info['texto'] for info in runs_info)
                            
                            # Traducir el párrafo completo
                            texto_traducido = traducir_texto(texto_completo, idioma_destino, idioma_origen)
                            
                            # Si la traducción falló, continuar con el siguiente párrafo
                            if texto_traducido is None:
                                print(f"Advertencia: Se mantuvo el texto original '{texto_completo}' debido a un error de traducción.")
                                continue
                            
                            # Distribuir la traducción entre los runs originales
                            if len(runs_info) == 1:
                                # Si solo hay un run, asignar toda la traducción a ese run
                                runs_info[0]['run'].text = texto_traducido
                            else:
                                # Si hay múltiples runs, intentar distribuir la traducción proporcionalmente
                                # Primero, asignar toda la traducción al primer run y vaciar los demás
                                runs_info[0]['run'].text = texto_traducido
                                for info in runs_info[1:]:
                                    info['run'].text = ""
                except Exception as e:
                    print(f"Error al procesar una forma en la diapositiva: {e}")
                    # Continuar con la siguiente forma
                    continue
                
                # Actualizar la barra de progreso
                barra_progreso.update(1)
        
        # Cerrar la barra de progreso
        barra_progreso.close()
        
        # Guardar la presentación traducida
        try:
            presentacion.save(ruta_archivo_traducido)
            print(f"\nPresentación traducida guardada como: {ruta_archivo_traducido}")
            return ruta_archivo_traducido
        except PermissionError:
            # Si no se puede guardar debido a permisos, intentar con un nombre alternativo
            import datetime
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            archivo_traducido = f"{nombre_sin_extension}_traducido_{idioma_destino}_{timestamp}.pptx"
            ruta_archivo_traducido = os.path.join(CARPETA_TRADUCIDOS, archivo_traducido)
            print(f"No se pudo guardar el archivo. Intentando con un nombre alternativo: {archivo_traducido}")
            try:
                presentacion.save(ruta_archivo_traducido)
                print(f"\nPresentación traducida guardada como: {ruta_archivo_traducido}")
                return ruta_archivo_traducido
            except Exception as e:
                print(f"Error al guardar la presentación con nombre alternativo: {e}")
                return None
    
    except Exception as e:
        print(f"Error al procesar la presentación: {e}")
        return None

def mover_a_originales(ruta_archivo):
    """
    Mueve un archivo a la carpeta de originales.
    
    Args:
        ruta_archivo (str): Ruta al archivo a mover
    
    Returns:
        str: Nueva ruta del archivo
    """
    # Verificar si el archivo ya está en la carpeta de originales
    if os.path.dirname(os.path.abspath(ruta_archivo)) == os.path.abspath(CARPETA_ORIGINALES):
        return ruta_archivo
    
    # Obtener el nombre del archivo
    nombre_archivo = os.path.basename(ruta_archivo)
    
    # Crear la nueva ruta
    nueva_ruta = os.path.join(CARPETA_ORIGINALES, nombre_archivo)
    
    # Copiar el archivo a la carpeta de originales
    shutil.copy2(ruta_archivo, nueva_ruta)
    
    print(f"Archivo copiado a: {nueva_ruta}")
    return nueva_ruta

def mostrar_ayuda():
    """Muestra instrucciones de uso del script."""
    print("\nTraductor de Presentaciones PowerPoint")
    print("=====================================")
    print("\nUso:")
    print("  python traductor_ppt.py <ruta_archivo> <idioma_destino> [idioma_origen]")
    print("\nEjemplos:")
    print("  python traductor_ppt.py presentacion.pptx es")
    print("  python traductor_ppt.py originales/presentacion.pptx fr es")
    print("\nCódigos de idioma comunes:")
    print("  es: Español")
    print("  en: Inglés")
    print("  fr: Francés")
    print("  de: Alemán")
    print("  it: Italiano")
    print("  pt: Portugués")
    print("  ja: Japonés")
    print("  zh-CN: Chino (Simplificado)")
    print("  ru: Ruso")
    print("\nNotas:")
    print("  - Los archivos originales se copiarán a la carpeta 'originales'")
    print("  - Los archivos traducidos se guardarán en la carpeta 'traducidos'")

def main():
    """Función principal del script."""
    # Verificar argumentos
    if len(sys.argv) < 3 or sys.argv[1] in ['-h', '--help']:
        mostrar_ayuda()
        return
    
    # Obtener argumentos
    ruta_archivo = sys.argv[1]
    idioma_destino = sys.argv[2]
    idioma_origen = sys.argv[3] if len(sys.argv) > 3 else 'auto'
    
    # Verificar que el archivo existe
    if not os.path.exists(ruta_archivo):
        print(f"Error: El archivo '{ruta_archivo}' no existe.")
        return
    
    # Verificar que el archivo es un PowerPoint
    if not ruta_archivo.lower().endswith(('.pptx', '.ppt')):
        print("Error: El archivo debe ser una presentación PowerPoint (.pptx o .ppt).")
        return
    
    # Mover el archivo a la carpeta de originales
    ruta_archivo = mover_a_originales(ruta_archivo)
    
    print(f"Traduciendo '{ruta_archivo}' de '{idioma_origen}' a '{idioma_destino}'...")
    traducir_presentacion(ruta_archivo, idioma_destino, idioma_origen)

if __name__ == "__main__":
    main() 
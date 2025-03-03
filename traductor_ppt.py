#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
Traductor de Presentaciones PowerPoint
Este script permite traducir automáticamente el contenido de presentaciones PowerPoint
de un idioma a otro utilizando la biblioteca deep-translator.
"""

import os
import sys
import time
import shutil
import re
from pptx import Presentation
from deep_translator import GoogleTranslator
from tqdm import tqdm
from dotenv import load_dotenv
import openai

# Cargar variables de entorno si existe un archivo .env
load_dotenv()

# Configurar la API de OpenAI
openai.api_key = os.getenv("OPENAI_API_KEY")

# Definir rutas de carpetas
CARPETA_ORIGINALES = "originales"
CARPETA_TRADUCIDOS = "traducidos"
CARPETA_COPIA_SEGURIDAD = "copia_seguridad"  # Añadimos la carpeta de copia de seguridad pero no la tocamos

# Crear carpetas si no existen
os.makedirs(CARPETA_ORIGINALES, exist_ok=True)
os.makedirs(CARPETA_TRADUCIDOS, exist_ok=True)

def traducir_texto_google(texto, idioma_destino='en', idioma_origen='auto'):
    """
    Traduce un texto usando Google Translate.
    
    Args:
        texto (str): Texto a traducir
        idioma_destino (str): Código del idioma de destino
        idioma_origen (str): Código del idioma de origen
    
    Returns:
        str: Texto traducido o el texto original si hay un error
    """
    try:
        # Añadir un pequeño retraso para evitar bloqueos de la API
        time.sleep(0.5)
        
        # Usar GoogleTranslator de deep_translator
        translator = GoogleTranslator(source=idioma_origen, target=idioma_destino)
        resultado = translator.translate(texto)
        
        # Verificar que el resultado no sea None
        if resultado is None:
            print(f"Advertencia: La traducción de '{texto}' devolvió None. Se usará el texto original.")
            return texto
            
        return resultado
    except Exception as e:
        print(f"Error al traducir el texto con Google Translate: {e}")
        return texto

def traducir_texto_openai(texto, idioma_destino='en', idioma_origen='auto'):
    """
    Traduce un texto usando la API de OpenAI.
    
    Args:
        texto (str): Texto a traducir
        idioma_destino (str): Código del idioma de destino
        idioma_origen (str): Código del idioma de origen
    
    Returns:
        str: Texto traducido o el texto original si hay un error
    """
    try:
        # Mapeo de códigos de idioma a nombres completos
        idiomas = {
            'en': 'English',
            'es': 'Spanish',
            'fr': 'French',
            'de': 'German',
            'it': 'Italian',
            'pt': 'Portuguese',
            'ru': 'Russian',
            'zh': 'Chinese',
            'ja': 'Japanese',
            'ko': 'Korean',
            'ar': 'Arabic',
            'hi': 'Hindi',
            'auto': 'auto-detected language'
        }
        
        # Obtener los nombres completos de los idiomas
        idioma_origen_nombre = idiomas.get(idioma_origen, idioma_origen)
        idioma_destino_nombre = idiomas.get(idioma_destino, idioma_destino)
        
        # Crear el mensaje para la API de OpenAI
        messages = [
            {"role": "system", "content": f"You are a professional translator. Translate the following text from {idioma_origen_nombre} to {idioma_destino_nombre}. Provide only the translation, no explanations or additional text."},
            {"role": "user", "content": texto}
        ]
        
        # Llamar a la API de OpenAI
        response = openai.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=messages,
            temperature=0.3,  # Valor bajo para obtener traducciones más consistentes
            max_tokens=1024
        )
        
        # Obtener la respuesta
        resultado = response.choices[0].message.content.strip()
        
        return resultado
    except Exception as e:
        print(f"Error al traducir el texto con OpenAI: {e}")
        return texto

def traducir_texto(texto, idioma_destino='en', idioma_origen='auto', metodo='google'):
    """
    Traduce un texto al idioma especificado usando el método indicado.
    
    Args:
        texto (str): Texto a traducir
        idioma_destino (str): Código del idioma de destino (por defecto 'en' para inglés)
        idioma_origen (str): Código del idioma de origen (por defecto 'auto' para detección automática)
        metodo (str): Método de traducción ('google' o 'openai')
    
    Returns:
        str: Texto traducido o el texto original si hay un error
    """
    if not texto or not texto.strip():
        return texto
    
    if metodo == 'openai':
        return traducir_texto_openai(texto, idioma_destino, idioma_origen)
    else:
        return traducir_texto_google(texto, idioma_destino, idioma_origen)

def traducir_presentacion(ruta_archivo, idioma_destino='en', idioma_origen='auto', metodo='google'):
    """
    Traduce una presentación PowerPoint al idioma especificado.
    
    Args:
        ruta_archivo (str): Ruta al archivo PowerPoint
        idioma_destino (str): Código del idioma de destino
        idioma_origen (str): Código del idioma de origen
        metodo (str): Método de traducción ('google' o 'openai')
    
    Returns:
        str: Ruta al archivo traducido
    """
    try:
        # Cargar la presentación
        presentacion = Presentation(ruta_archivo)
        
        # Obtener el nombre del archivo sin extensión
        nombre_base = os.path.basename(ruta_archivo)
        nombre_sin_extension = os.path.splitext(nombre_base)[0]
        
        # Crear el nombre del archivo traducido
        archivo_traducido = f"{nombre_sin_extension}_traducido_{idioma_destino}.pptx"
        ruta_archivo_traducido = os.path.join(CARPETA_TRADUCIDOS, archivo_traducido)
        
        # Verificar si el archivo ya existe y está en uso
        if os.path.exists(ruta_archivo_traducido):
            try:
                # Intentar abrir el archivo para verificar si está en uso
                with open(ruta_archivo_traducido, 'a'):
                    pass
            except PermissionError:
                # Si está en uso, crear un nombre alternativo con timestamp
                import datetime
                timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
                archivo_traducido = f"{nombre_sin_extension}_traducido_{idioma_destino}_{timestamp}.pptx"
                ruta_archivo_traducido = os.path.join(CARPETA_TRADUCIDOS, archivo_traducido)
                print(f"El archivo original está en uso. Se usará un nombre alternativo: {archivo_traducido}")
        
        # Contar el número total de elementos a traducir para la barra de progreso
        total_elementos = sum(1 for diapositiva in presentacion.slides for forma in diapositiva.shapes if hasattr(forma, "text_frame"))
        
        # Inicializar la barra de progreso
        barra_progreso = tqdm(total=total_elementos, desc="Traduciendo", unit="elementos")
        
        # Procesar cada diapositiva
        for diapositiva in presentacion.slides:
            # Procesar cada forma en la diapositiva
            for forma in diapositiva.shapes:
                # Verificar si la forma tiene texto
                if hasattr(forma, "text_frame"):
                    try:
                        # Procesar cada párrafo en el marco de texto
                        for parrafo in forma.text_frame.paragraphs:
                            # Si el párrafo está vacío, continuar con el siguiente
                            if not parrafo.text.strip():
                                continue
                            
                            # Recopilar información sobre cada run en el párrafo
                            runs_info = []
                            texto_completo = ""
                            
                            for run in parrafo.runs:
                                # Si el run está vacío, continuar con el siguiente
                                if not run.text.strip():
                                    continue
                                
                                # Añadir el texto del run al texto completo
                                texto_completo += run.text
                                
                                # Guardar información sobre el run
                                runs_info.append({
                                    'run': run,
                                    'texto_original': run.text
                                })
                            
                            # Si no hay texto que traducir, continuar con el siguiente párrafo
                            if not texto_completo.strip():
                                continue
                            
                            # Traducir el párrafo completo
                            texto_traducido = traducir_texto(texto_completo, idioma_destino, idioma_origen, metodo)
                            
                            # Si la traducción falló, continuar con el siguiente párrafo
                            if texto_traducido is None:
                                print(f"Advertencia: Se mantuvo el texto original '{texto_completo}' debido a un error de traducción.")
                                continue
                            
                            # Distribuir la traducción entre los runs originales
                            if len(runs_info) == 1:
                                # Si solo hay un run, asignar toda la traducción a ese run
                                runs_info[0]['run'].text = texto_traducido
                            else:
                                # Si hay múltiples runs, intentar distribuir la traducción proporcionalmente
                                # Primero, asignar toda la traducción al primer run y vaciar los demás
                                runs_info[0]['run'].text = texto_traducido
                                for info in runs_info[1:]:
                                    info['run'].text = ""
                    except Exception as e:
                        print(f"Error al procesar una forma en la diapositiva: {e}")
                        # Continuar con la siguiente forma
                        continue
                
                # Actualizar la barra de progreso
                barra_progreso.update(1)
        
        # Cerrar la barra de progreso
        barra_progreso.close()
        
        # Guardar la presentación traducida
        try:
            presentacion.save(ruta_archivo_traducido)
            print(f"\nPresentación traducida guardada como: {ruta_archivo_traducido}")
            return ruta_archivo_traducido
        except PermissionError:
            # Si no se puede guardar debido a permisos, intentar con un nombre alternativo
            import datetime
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            archivo_traducido = f"{nombre_sin_extension}_traducido_{idioma_destino}_{timestamp}.pptx"
            ruta_archivo_traducido = os.path.join(CARPETA_TRADUCIDOS, archivo_traducido)
            print(f"No se pudo guardar el archivo. Intentando con un nombre alternativo: {archivo_traducido}")
            try:
                presentacion.save(ruta_archivo_traducido)
                print(f"\nPresentación traducida guardada como: {ruta_archivo_traducido}")
                return ruta_archivo_traducido
            except Exception as e:
                print(f"Error al guardar la presentación con nombre alternativo: {e}")
                return None
    
    except Exception as e:
        print(f"Error al procesar la presentación: {e}")
        return None

def mover_a_originales(ruta_archivo):
    """
    Mueve un archivo a la carpeta de originales.
    
    Args:
        ruta_archivo (str): Ruta al archivo a mover
    
    Returns:
        str: Nueva ruta del archivo
    """
    # Verificar si el archivo ya está en la carpeta de originales
    if os.path.dirname(os.path.abspath(ruta_archivo)) == os.path.abspath(CARPETA_ORIGINALES):
        return ruta_archivo
    
    # Obtener el nombre del archivo
    nombre_archivo = os.path.basename(ruta_archivo)
    
    # Crear la nueva ruta
    nueva_ruta = os.path.join(CARPETA_ORIGINALES, nombre_archivo)
    
    # Copiar el archivo a la carpeta de originales
    shutil.copy2(ruta_archivo, nueva_ruta)
    
    print(f"Archivo copiado a: {nueva_ruta}")
    return nueva_ruta

def mostrar_ayuda():
    """Muestra la ayuda del script."""
    print("""
Uso: python traductor_ppt.py <archivo_powerpoint> <idioma_destino> [idioma_origen] [metodo_traduccion]

Argumentos:
  archivo_powerpoint   Ruta al archivo PowerPoint a traducir
  idioma_destino       Código del idioma de destino (ej: 'en' para inglés)
  idioma_origen        (Opcional) Código del idioma de origen (por defecto 'auto' para detección automática)
  metodo_traduccion    (Opcional) Método de traducción ('google' o 'openai', por defecto 'google')

Ejemplos:
  python traductor_ppt.py presentacion.pptx en es
  python traductor_ppt.py presentacion.pptx fr auto
  python traductor_ppt.py presentacion.pptx en es openai

Nota: Para usar OpenAI, debes tener una clave de API válida en un archivo .env con la variable OPENAI_API_KEY
""")

def main():
    """Función principal del script."""
    # Verificar argumentos
    if len(sys.argv) < 3 or sys.argv[1] in ['-h', '--help']:
        mostrar_ayuda()
        return
    
    # Obtener argumentos
    ruta_archivo = sys.argv[1]
    idioma_destino = sys.argv[2]
    idioma_origen = sys.argv[3] if len(sys.argv) > 3 else 'auto'
    metodo_traduccion = sys.argv[4] if len(sys.argv) > 4 else 'google'
    
    # Verificar que el método de traducción es válido
    if metodo_traduccion not in ['google', 'openai']:
        print(f"Error: El método de traducción '{metodo_traduccion}' no es válido. Debe ser 'google' o 'openai'.")
        return
    
    # Verificar que el archivo existe
    if not os.path.exists(ruta_archivo):
        print(f"Error: El archivo '{ruta_archivo}' no existe.")
        return
    
    # Verificar que el archivo es un PowerPoint
    if not ruta_archivo.lower().endswith(('.pptx', '.ppt')):
        print("Error: El archivo debe ser una presentación PowerPoint (.pptx o .ppt).")
        return
    
    # Si se usa OpenAI, verificar que existe la clave de API
    if metodo_traduccion == 'openai' and not os.getenv("OPENAI_API_KEY"):
        print("Error: Para usar OpenAI, debes tener una clave de API válida en un archivo .env con la variable OPENAI_API_KEY.")
        return
    
    # Mover el archivo a la carpeta de originales
    ruta_archivo = mover_a_originales(ruta_archivo)
    
    print(f"Traduciendo '{ruta_archivo}' de '{idioma_origen}' a '{idioma_destino}' usando el método '{metodo_traduccion}'...")
    traducir_presentacion(ruta_archivo, idioma_destino, idioma_origen, metodo_traduccion)

if __name__ == "__main__":
    main() 